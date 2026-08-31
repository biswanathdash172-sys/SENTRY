"""
routers/reports.py
-------------------
Enhanced PPT export — professional SENTRY SOC branded template.

Slides produced:
  1.  Cover       — dark gradient background, SENTRY logo text, org ID, timestamp
  2.  TOC         — table of contents card layout
  3.  Executive   — 4 KPI "stat card" boxes (total scanned, high risk %, resolved %, avg time)
  4.  Risk Analysis Summary — structured risk tier & resolution breakdown with colour coding
  5.  Tier Chart  — native bar chart, coloured bars per tier
  6.  Trend Chart — native line-markers chart for last 7 days
  7.  Resolution Breakdown — stacked horizontal bar chart (auto / admin / denied)
  8.  Risk Analysis Deep Dive — table of tier × status cross-tab
  9.  Footer / Disclaimer slide
"""

from __future__ import annotations

import io
from datetime import datetime, timezone

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from models import AdminUser
from routers.org_auth import require_admin
from routers.analytics import get_analytics

router = APIRouter(tags=["reports"])

# ── Brand palette (matches dashboard CSS variables) ──────────────────────────
C_BG         = RGBColor(0x0A, 0x0E, 0x27)   # --bg
C_PANEL      = RGBColor(0x10, 0x14, 0x2E)   # --panel
C_ACCENT     = RGBColor(0x3B, 0x5B, 0xFF)   # --accent  (blue)
C_CYAN       = RGBColor(0x22, 0xD3, 0xEE)   # --accent-2 (cyan)
C_NOT_RISKY  = RGBColor(0x22, 0xD3, 0xEE)   # low risk
C_PART_RISKY = RGBColor(0xF5, 0x9E, 0x0B)   # medium risk
C_HIGH_RISKY = RGBColor(0xEF, 0x5B, 0x5B)   # high risk
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED      = RGBColor(0x94, 0xA3, 0xB8)
C_SUCCESS    = RGBColor(0x22, 0xC5, 0x5E)
C_BORDER     = RGBColor(0x1E, 0x26, 0x4A)   # slightly lighter than panel

SLIDE_W = Inches(13.33)   # 16:9 widescreen
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _solid_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill: RGBColor, alpha: int = 255):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()   # no border
    return shape


def _textbox(slide, x, y, w, h, text, size=Pt(18), bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _accent_bar(slide, color=C_CYAN, height=Inches(0.06)):
    """Thin accent stripe across full slide width."""
    _rect(slide, 0, 0, SLIDE_W, height, color)


def _slide_footer(slide, page_num: int, total: int, org_id: str):
    """Bottom strip with org name, page number."""
    bar_h = Inches(0.35)
    _rect(slide, 0, SLIDE_H - bar_h, SLIDE_W, bar_h, C_PANEL)
    _textbox(slide,
             Inches(0.4), SLIDE_H - bar_h + Inches(0.05),
             Inches(6), bar_h,
             f"SENTRY SOC  ·  {org_id}  ·  Confidential",
             size=Pt(9), color=C_MUTED)
    _textbox(slide,
             SLIDE_W - Inches(2), SLIDE_H - bar_h + Inches(0.05),
             Inches(1.8), bar_h,
             f"{page_num} / {total}",
             size=Pt(9), color=C_MUTED, align=PP_ALIGN.RIGHT)


def _section_title(slide, title: str, y=Inches(0.55)):
    """Section heading with cyan underline accent."""
    _textbox(slide, Inches(0.55), y, Inches(10), Inches(0.55),
             title.upper(), size=Pt(22), bold=True, color=C_CYAN)
    # Underline rectangle
    _rect(slide, Inches(0.55), y + Inches(0.5), Inches(2.5), Inches(0.025), C_ACCENT)


def _kpi_card(slide, x, y, w, h, label, value, value_color=C_WHITE, sub=""):
    """A single KPI stat card box."""
    _rect(slide, x, y, w, h, C_PANEL)
    # subtle left accent strip
    _rect(slide, x, y, Inches(0.05), h, C_ACCENT)
    _textbox(slide, x + Inches(0.15), y + Inches(0.12), w - Inches(0.2), Inches(0.3),
             label, size=Pt(9), color=C_MUTED)
    _textbox(slide, x + Inches(0.15), y + Inches(0.38), w - Inches(0.2), Inches(0.55),
             str(value), size=Pt(30), bold=True, color=value_color)
    if sub:
        _textbox(slide, x + Inches(0.15), y + h - Inches(0.28), w - Inches(0.2), Inches(0.24),
                 sub, size=Pt(8), color=C_MUTED)


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_cover(prs, org_id: str, generated_at: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)

    # Left accent gradient panel
    _rect(slide, 0, 0, Inches(0.45), SLIDE_H, C_ACCENT)
    _rect(slide, Inches(0.45), 0, Inches(0.15), SLIDE_H, C_CYAN)

    # Top thin accent bar
    _rect(slide, Inches(0.6), 0, SLIDE_W, Inches(0.04), C_CYAN)

    # Decorative hex-like polygon (simulated with large rectangle + overlay)
    _rect(slide, SLIDE_W - Inches(4.5), Inches(0.5), Inches(4), Inches(6),
          RGBColor(0x0D, 0x12, 0x32))
    _rect(slide, SLIDE_W - Inches(4.4), Inches(0.6), Inches(3.8), Inches(5.8),
          RGBColor(0x10, 0x14, 0x2E))
    # Cyan ring simulation
    _rect(slide, SLIDE_W - Inches(3.2), Inches(2.2), Inches(1.5), Inches(1.5),
          RGBColor(0x05, 0x28, 0x40))

    # SENTRY wordmark
    _textbox(slide, Inches(0.9), Inches(1.5), Inches(7), Inches(1.2),
             "SENTRY", size=Pt(64), bold=True, color=C_WHITE)
    _textbox(slide, Inches(0.9), Inches(2.55), Inches(7), Inches(0.5),
             "Security Operations Centre  ·  Risk Intelligence Report",
             size=Pt(14), color=C_CYAN)

    # Divider line
    _rect(slide, Inches(0.9), Inches(3.2), Inches(5.5), Inches(0.025), C_ACCENT)

    # Org + date
    _textbox(slide, Inches(0.9), Inches(3.4), Inches(7), Inches(0.45),
             f"Organisation:   {org_id}", size=Pt(14), bold=True, color=C_WHITE)
    _textbox(slide, Inches(0.9), Inches(3.85), Inches(7), Inches(0.35),
             f"Generated:       {generated_at}", size=Pt(12), color=C_MUTED)
    _textbox(slide, Inches(0.9), Inches(4.25), Inches(7), Inches(0.35),
             "Classification:  CONFIDENTIAL — Internal Use Only",
             size=Pt(11), color=RGBColor(0xEF, 0x5B, 0x5B))

    # Bottom tag
    _textbox(slide, Inches(0.9), SLIDE_H - Inches(0.8), Inches(5), Inches(0.35),
             "Powered by SENTRY SOC Platform",
             size=Pt(10), color=C_MUTED)


def _slide_toc(prs, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "Table of Contents")
    _slide_footer(slide, 2, 9, org_id)

    sections = [
        ("01", "Executive Summary", "Key metrics at a glance"),
        ("02", "Risk Analysis Summary", "Tier breakdown and resolution status"),
        ("03", "Risk Distribution Chart", "Visual breakdown by risk tier"),
        ("04", "7-Day Risk Trend", "Daily frequency over the last week"),
        ("05", "Resolution Breakdown", "How findings were resolved"),
        ("06", "Risk Deep Dive", "Tier vs. status cross-analysis table"),
        ("07", "Disclaimer", "Report scope and limitations"),
    ]

    for i, (num, title, sub) in enumerate(sections):
        row = i % 4
        col = i // 4
        x = Inches(0.6) + col * Inches(6.5)
        y = Inches(1.4) + row * Inches(1.35)
        w, h = Inches(5.8), Inches(1.1)
        _rect(slide, x, y, w, h, C_PANEL)
        _rect(slide, x, y, Inches(0.06), h, C_ACCENT)
        _textbox(slide, x + Inches(0.18), y + Inches(0.08), Inches(0.5), Inches(0.35),
                 num, size=Pt(11), bold=True, color=C_CYAN)
        _textbox(slide, x + Inches(0.65), y + Inches(0.07), w - Inches(0.75), Inches(0.35),
                 title, size=Pt(13), bold=True, color=C_WHITE)
        _textbox(slide, x + Inches(0.65), y + Inches(0.48), w - Inches(0.75), Inches(0.35),
                 sub, size=Pt(10), color=C_MUTED)


def _slide_executive(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "01 · Executive Summary")
    _slide_footer(slide, 3, 9, org_id)

    tc = analytics.tier_counts
    sc = analytics.status_counts
    total = analytics.total_scanned or 1
    high_pct = f"{tc['high_risky'] / total * 100:.1f}%"
    res_pct  = f"{sc['completed'] / total * 100:.1f}%"
    avg_time = f"{analytics.avg_resolution_minutes} min" if analytics.avg_resolution_minutes else "N/A"

    card_w, card_h = Inches(2.8), Inches(1.35)
    gap = Inches(0.28)
    y = Inches(1.4)
    cards = [
        ("Total Scanned",    analytics.total_scanned, C_CYAN,       "risk items found"),
        ("High Risk Items",  tc["high_risky"],         C_HIGH_RISKY, f"{high_pct} of total"),
        ("Resolved",         sc["completed"],          C_SUCCESS,    f"{res_pct} resolution rate"),
        ("Avg. Resolution",  avg_time,                 C_ACCENT,     "time per finding"),
    ]
    for i, (label, val, color, sub) in enumerate(cards):
        _kpi_card(slide,
                  Inches(0.55) + i * (card_w + gap), y,
                  card_w, card_h, label, val, color, sub)

    # Second row — pending, part risky, not risky, auto-approved
    rb = analytics.resolution_breakdown
    y2 = y + card_h + Inches(0.35)
    cards2 = [
        ("Pending Review",   sc["pending"],               C_PART_RISKY, "awaiting action"),
        ("Part Risky",       tc["part_risky"],            C_PART_RISKY, "medium severity"),
        ("Not Risky",        tc["not_risky"],             C_NOT_RISKY,  "low severity"),
        ("Auto-Approved",    rb["auto_approved"],         C_MUTED,      "by policy rules"),
    ]
    for i, (label, val, color, sub) in enumerate(cards2):
        _kpi_card(slide,
                  Inches(0.55) + i * (card_w + gap), y2,
                  card_w, card_h, label, val, color, sub)

    # Insight callout
    _rect(slide, Inches(0.55), SLIDE_H - Inches(1.45), SLIDE_W - Inches(1.1), Inches(0.9), C_PANEL)
    insight = (
        f"  [!] {tc['high_risky']} HIGH-RISK item(s) require immediate attention.   "
        f"  [OK] {sc['completed']} item(s) have been resolved.   "
        f"  Avg. resolution time: {avg_time}."
    )
    _textbox(slide, Inches(0.65), SLIDE_H - Inches(1.4), SLIDE_W - Inches(1.3), Inches(0.85),
             insight, size=Pt(11), color=C_WHITE, wrap=True)


def _slide_risk_summary(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "02 · Risk Analysis Summary")
    _slide_footer(slide, 4, 9, org_id)

    tc = analytics.tier_counts
    sc = analytics.status_counts
    rb = analytics.resolution_breakdown
    total = analytics.total_scanned or 1

    # Left column — tier breakdown
    _rect(slide, Inches(0.55), Inches(1.45), Inches(5.5), Inches(5.3), C_PANEL)
    _textbox(slide, Inches(0.7), Inches(1.55), Inches(5.2), Inches(0.4),
             "Risk Tier Breakdown", size=Pt(13), bold=True, color=C_CYAN)

    tiers = [
        ("Not Risky",   tc["not_risky"],   C_NOT_RISKY),
        ("Part Risky",  tc["part_risky"],  C_PART_RISKY),
        ("High Risk",   tc["high_risky"],  C_HIGH_RISKY),
    ]
    for i, (label, count, color) in enumerate(tiers):
        ty = Inches(2.1) + i * Inches(1.35)
        pct = count / total * 100
        # Bar background
        _rect(slide, Inches(0.7), ty, Inches(4.8), Inches(0.3), C_BORDER)
        # Filled portion
        bar_w = max(Inches(0.05), Inches(4.8) * pct / 100)
        _rect(slide, Inches(0.7), ty, bar_w, Inches(0.3), color)
        _textbox(slide, Inches(0.7), ty - Inches(0.28), Inches(3), Inches(0.28),
                 label, size=Pt(11), bold=True, color=C_WHITE)
        _textbox(slide, Inches(4.2), ty - Inches(0.28), Inches(1.2), Inches(0.28),
                 f"{count}  ({pct:.1f}%)", size=Pt(11), color=color, align=PP_ALIGN.RIGHT)

    # Right column — resolution & method
    _rect(slide, Inches(6.55), Inches(1.45), Inches(5.9), Inches(2.35), C_PANEL)
    _textbox(slide, Inches(6.7), Inches(1.55), Inches(5.5), Inches(0.4),
             "Resolution Status", size=Pt(13), bold=True, color=C_CYAN)
    statuses = [
        ("Pending",   sc["pending"],   C_PART_RISKY),
        ("Completed", sc["completed"], C_SUCCESS),
    ]
    for i, (label, count, color) in enumerate(statuses):
        sy = Inches(2.1) + i * Inches(0.75)
        _textbox(slide, Inches(6.7), sy, Inches(3), Inches(0.45),
                 f"{label}:", size=Pt(12), color=C_MUTED)
        _textbox(slide, Inches(9.5), sy, Inches(2.5), Inches(0.45),
                 str(count), size=Pt(20), bold=True, color=color, align=PP_ALIGN.RIGHT)

    _rect(slide, Inches(6.55), Inches(3.9), Inches(5.9), Inches(2.85), C_PANEL)
    _textbox(slide, Inches(6.7), Inches(4.0), Inches(5.5), Inches(0.4),
             "Resolution Method", size=Pt(13), bold=True, color=C_CYAN)
    methods = [
        ("Auto-Approved",  rb["auto_approved"],  C_MUTED),
        ("Admin-Approved", rb["admin_approved"], C_SUCCESS),
        ("Admin-Denied",   rb["admin_denied"],   C_HIGH_RISKY),
    ]
    for i, (label, count, color) in enumerate(methods):
        my = Inches(4.55) + i * Inches(0.72)
        _textbox(slide, Inches(6.7), my, Inches(3), Inches(0.45),
                 f"{label}:", size=Pt(11), color=C_MUTED)
        _textbox(slide, Inches(9.5), my, Inches(2.5), Inches(0.45),
                 str(count), size=Pt(18), bold=True, color=color, align=PP_ALIGN.RIGHT)

    # Avg resolution time callout
    avg_time = f"{analytics.avg_resolution_minutes} min" if analytics.avg_resolution_minutes else "No data yet"
    _rect(slide, Inches(0.55), SLIDE_H - Inches(0.95), SLIDE_W - Inches(1.1), Inches(0.55), C_PANEL)
    _rect(slide, Inches(0.55), SLIDE_H - Inches(0.95), Inches(0.06), Inches(0.55), C_CYAN)
    _textbox(slide, Inches(0.75), SLIDE_H - Inches(0.93), Inches(8), Inches(0.5),
             f"Average Resolution Time:  {avg_time}", size=Pt(13), bold=True, color=C_WHITE)


def _slide_tier_chart(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "03 · Risk Distribution by Tier")
    _slide_footer(slide, 5, 9, org_id)

    tc = analytics.tier_counts
    chart_data = CategoryChartData()
    chart_data.categories = ["Not Risky", "Part Risky", "High Risky"]
    chart_data.add_series("Items", (
        tc["not_risky"], tc["part_risky"], tc["high_risky"]
    ))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6), Inches(1.35), Inches(12), Inches(5.45), chart_data
    )
    chart = gf.chart
    chart.has_legend = False
    chart.plots[0].has_data_labels = True

    series = chart.series[0]
    tier_colors = [C_NOT_RISKY, C_PART_RISKY, C_HIGH_RISKY]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = tier_colors[i]

    # Style chart background
    chart.chart_area.fill.solid()
    chart.chart_area.fill.fore_color.rgb = C_PANEL
    chart.plot_area.fill.solid()
    chart.plot_area.fill.fore_color.rgb = C_BG


def _slide_trend_chart(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "04 · 7-Day Risk Frequency Trend")
    _slide_footer(slide, 6, 9, org_id)

    dc = analytics.daily_counts
    if not dc:
        _textbox(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1),
                 "No trend data available for the selected period.",
                 size=Pt(16), color=C_MUTED, align=PP_ALIGN.CENTER)
        return

    chart_data = CategoryChartData()
    chart_data.categories = [d["date"] for d in dc]
    chart_data.add_series("Risks Detected", tuple(d["count"] for d in dc))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.6), Inches(1.35), Inches(12), Inches(5.45), chart_data
    )
    chart = gf.chart
    chart.has_legend = False

    series = chart.series[0]
    series.format.line.color.rgb = C_ACCENT
    series.format.line.width = Pt(2.5)

    chart.chart_area.fill.solid()
    chart.chart_area.fill.fore_color.rgb = C_PANEL
    chart.plot_area.fill.solid()
    chart.plot_area.fill.fore_color.rgb = C_BG


def _slide_resolution_chart(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "05 · Resolution Method Breakdown")
    _slide_footer(slide, 7, 9, org_id)

    rb = analytics.resolution_breakdown
    chart_data = CategoryChartData()
    chart_data.categories = ["Resolution Methods"]
    chart_data.add_series("Auto-Approved",  (rb["auto_approved"],))
    chart_data.add_series("Admin-Approved", (rb["admin_approved"],))
    chart_data.add_series("Admin-Denied",   (rb["admin_denied"],))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.6), Inches(1.35), Inches(12), Inches(4.5), chart_data
    )
    chart = gf.chart
    chart.has_legend = True

    colors = [C_MUTED, C_SUCCESS, C_HIGH_RISKY]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i]

    chart.chart_area.fill.solid()
    chart.chart_area.fill.fore_color.rgb = C_PANEL
    chart.plot_area.fill.solid()
    chart.plot_area.fill.fore_color.rgb = C_BG

    # Summary row below chart
    _rect(slide, Inches(0.6), Inches(6.05), Inches(12), Inches(0.8), C_PANEL)
    rb_text = (
        f"  Auto-Approved: {rb['auto_approved']}    "
        f"Admin-Approved: {rb['admin_approved']}    "
        f"Admin-Denied: {rb['admin_denied']}    "
        f"Total Resolved: {analytics.status_counts['completed']}"
    )
    _textbox(slide, Inches(0.75), Inches(6.1), Inches(11.5), Inches(0.7),
             rb_text, size=Pt(12), color=C_WHITE)


def _slide_deep_dive(prs, analytics, org_id: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "06 · Risk Deep Dive — Tier × Status")
    _slide_footer(slide, 8, 9, org_id)

    tc = analytics.tier_counts
    sc = analytics.status_counts
    rb = analytics.resolution_breakdown
    total = analytics.total_scanned or 1

    # Main table
    rows = 5   # header + 3 tiers + totals
    cols = 5   # Tier | Count | % of Total | Pending | Resolved
    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.45), Inches(12), Inches(4.2)
    ).table

    headers = ["Risk Tier", "Total Items", "% of Scanned", "Pending", "Resolved"]
    header_colors = [C_ACCENT, C_ACCENT, C_ACCENT, C_ACCENT, C_ACCENT]
    for c, (h, hc) in enumerate(zip(headers, header_colors)):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hc
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = C_WHITE

    data_rows = [
        ("Not Risky",  tc["not_risky"],  C_NOT_RISKY),
        ("Part Risky", tc["part_risky"], C_PART_RISKY),
        ("High Risk",  tc["high_risky"], C_HIGH_RISKY),
    ]
    # Approximate pending/resolved split per tier (proportional estimate)
    for r, (tier_name, count, color) in enumerate(data_rows, start=1):
        pct = count / total * 100
        # Estimate: use same overall ratio for pending / completed
        est_pending = round(count * sc["pending"] / total) if total else 0
        est_resolved = count - est_pending
        row_data = [tier_name, str(count), f"{pct:.1f}%", str(est_pending), str(est_resolved)]
        row_bg = C_PANEL if r % 2 == 0 else C_BG
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(12)
            run.font.color.rgb = color if c == 0 else C_WHITE
            run.font.bold = (c == 0)

    # Totals row
    totals = ["TOTAL", str(total), "100%", str(sc["pending"]), str(sc["completed"])]
    for c, val in enumerate(totals):
        cell = tbl.cell(4, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_ACCENT
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = C_WHITE

    # Note about estimation
    _textbox(slide, Inches(0.6), Inches(5.75), Inches(12), Inches(0.4),
             "* Pending / Resolved split per tier is a proportional estimate based on overall resolution rate.",
             size=Pt(9), color=C_MUTED)

    # Resolution method summary boxes
    meth_y = Inches(6.25)
    methods = [
        ("Auto-Approved", rb["auto_approved"], C_MUTED),
        ("Admin-Approved", rb["admin_approved"], C_SUCCESS),
        ("Admin-Denied", rb["admin_denied"], C_HIGH_RISKY),
    ]
    box_w = Inches(3.8)
    for i, (label, val, color) in enumerate(methods):
        bx = Inches(0.6) + i * (box_w + Inches(0.2))
        _rect(slide, bx, meth_y, box_w, Inches(0.65), C_PANEL)
        _rect(slide, bx, meth_y, Inches(0.05), Inches(0.65), color)
        _textbox(slide, bx + Inches(0.15), meth_y + Inches(0.05), Inches(2), Inches(0.3),
                 label, size=Pt(10), color=C_MUTED)
        _textbox(slide, bx + Inches(0.15), meth_y + Inches(0.3), Inches(1.5), Inches(0.32),
                 str(val), size=Pt(16), bold=True, color=color)


def _slide_disclaimer(prs, org_id: str, generated_at: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, C_BG)
    _accent_bar(slide)
    _section_title(slide, "07 · Disclaimer & Report Scope")
    _slide_footer(slide, 9, 9, org_id)

    _rect(slide, Inches(0.6), Inches(1.45), SLIDE_W - Inches(1.2), Inches(4.8), C_PANEL)
    _rect(slide, Inches(0.6), Inches(1.45), Inches(0.07), Inches(4.8), C_HIGH_RISKY)

    disclaimer_lines = [
        "1.  This report is generated automatically from live Supabase data at the time of export.",
        "2.  Risk classifications (Not Risky / Part Risky / High Risk) are based on the SENTRY SOC",
        "     scanning engine and may not reflect all real-world threat vectors.",
        "3.  Resolution estimates per tier are proportional approximations unless per-finding",
        "     breakdowns are explicitly tracked in the database.",
        "4.  This document is CONFIDENTIAL and intended for internal use by authorised personnel only.",
        "5.  The SENTRY platform does not provide legal, compliance, or regulatory assurance.",
        "     Independent audits are recommended for regulatory filings.",
        "6.  For questions about this report, contact your SENTRY SOC administrator.",
    ]

    for i, line in enumerate(disclaimer_lines):
        _textbox(slide, Inches(0.85), Inches(1.65) + i * Inches(0.47),
                 SLIDE_W - Inches(1.5), Inches(0.45),
                 line, size=Pt(11), color=C_WHITE if i == 0 else C_MUTED, wrap=True)

    _textbox(slide, Inches(0.6), SLIDE_H - Inches(1.4), SLIDE_W - Inches(1.2), Inches(0.4),
             f"Report generated: {generated_at}  ·  Organisation: {org_id}",
             size=Pt(10), color=C_MUTED)
    _textbox(slide, Inches(0.6), SLIDE_H - Inches(1.0), SLIDE_W - Inches(1.2), Inches(0.4),
             "© SENTRY SOC Platform — All rights reserved.",
             size=Pt(10), color=C_MUTED)


# ── Main export route ─────────────────────────────────────────────────────────

@router.get("/reports/export")
def export_report(admin: AdminUser = Depends(require_admin)):
    """
    ADMIN ONLY. Generates a fully-designed .pptx report from live data
    and streams it back as a file download. Nothing is saved to disk.
    """
    analytics = get_analytics(admin=admin)
    generated_at = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, admin.org_id, generated_at)
    _slide_toc(prs, admin.org_id)
    _slide_executive(prs, analytics, admin.org_id)
    _slide_risk_summary(prs, analytics, admin.org_id)
    _slide_tier_chart(prs, analytics, admin.org_id)
    _slide_trend_chart(prs, analytics, admin.org_id)
    _slide_resolution_chart(prs, analytics, admin.org_id)
    _slide_deep_dive(prs, analytics, admin.org_id)
    _slide_disclaimer(prs, admin.org_id, generated_at)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = f"SENTRY-RiskReport-{admin.org_id}-{datetime.now(timezone.utc).strftime('%Y%m%d-%H%M%S')}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )